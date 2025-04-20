import pdfplumber
from pptx import Presentation
from fastapi import UploadFile

async def parse_file(file: UploadFile) -> str:
    content = ""
    if file.filename.endswith(".pdf"):
        with pdfplumber.open(file.file) as pdf:
            for page in pdf.pages:
                content += page.extract_text() + "\n"
    elif file.filename.endswith(".pptx"):
        prs = Presentation(file.file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
    else:
        content = (await file.read()).decode("utf-8")
    return content
